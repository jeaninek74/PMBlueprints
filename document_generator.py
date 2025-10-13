"""
Document Generator for PMBlueprints
Generates professional Word, Excel, and PowerPoint documents
Integrates with PM intelligence for methodology-aware content
"""

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import os
import tempfile
from datetime import datetime
import logging

logger = logging.getLogger(__name__)


class DocumentGenerator:
    """Generate professional PM documents in various formats"""
    
    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
    
    def generate_word_document(self, content, document_name, structure, metadata):
        """
        Generate a professional Word document
        
        Args:
            content: Generated content from AI
            document_name: Name of the document
            structure: Document structure from PM intelligence
            metadata: Document metadata (methodology, PMBOK area, etc.)
        
        Returns:
            Path to generated Word file
        """
        try:
            # Create document
            doc = Document()
            
            # Set up styles
            self._setup_word_styles(doc)
            
            # Add title page
            self._add_word_title_page(doc, document_name, metadata)
            
            # Add page break
            doc.add_page_break()
            
            # Add table of contents placeholder
            doc.add_heading('Table of Contents', level=1)
            doc.add_paragraph('[Table of Contents will be generated in Microsoft Word]')
            doc.add_page_break()
            
            # Parse and add content
            self._add_word_content(doc, content, structure)
            
            # Add footer
            self._add_word_footer(doc, metadata)
            
            # Save document
            filename = f"{document_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            filepath = os.path.join(self.temp_dir, filename)
            doc.save(filepath)
            
            logger.info(f"Word document generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating Word document: {e}")
            raise
    
    def generate_excel_document(self, content, document_name, structure, metadata):
        """
        Generate a professional Excel spreadsheet
        
        Args:
            content: Generated content from AI
            document_name: Name of the document
            structure: Document structure from PM intelligence
            metadata: Document metadata
        
        Returns:
            Path to generated Excel file
        """
        try:
            # Create workbook
            wb = Workbook()
            ws = wb.active
            ws.title = document_name[:31]  # Excel sheet name limit
            
            # Add title
            ws['A1'] = document_name
            ws['A1'].font = Font(size=16, bold=True, color="FFFFFF")
            ws['A1'].fill = PatternFill(start_color="0066CC", end_color="0066CC", fill_type="solid")
            ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
            ws.merge_cells('A1:F1')
            ws.row_dimensions[1].height = 30
            
            # Add metadata
            row = 2
            ws[f'A{row}'] = f"Methodology: {metadata.get('methodology', 'N/A')}"
            ws[f'A{row}'].font = Font(italic=True)
            row += 1
            ws[f'A{row}'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            ws[f'A{row}'].font = Font(italic=True)
            row += 2
            
            # Add headers
            columns = structure.get('sheets', [{}])[0].get('columns', [
                'ID', 'Description', 'Owner', 'Status', 'Priority', 'Notes'
            ])
            
            for col_idx, column_name in enumerate(columns, start=1):
                cell = ws.cell(row=row, column=col_idx, value=column_name)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="0066CC", end_color="0066CC", fill_type="solid")
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
            
            # Add sample data rows
            self._add_excel_sample_data(ws, row + 1, columns, content, metadata)
            
            # Auto-adjust column widths
            for col_idx in range(1, len(columns) + 1):
                ws.column_dimensions[get_column_letter(col_idx)].width = 20
            
            # Freeze header row
            ws.freeze_panes = f'A{row + 1}'
            
            # Add auto-filter
            ws.auto_filter.ref = f'A{row}:{get_column_letter(len(columns))}{row + 10}'
            
            # Save workbook
            filename = f"{document_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            filepath = os.path.join(self.temp_dir, filename)
            wb.save(filepath)
            
            logger.info(f"Excel document generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating Excel document: {e}")
            raise
    
    def generate_powerpoint_document(self, content, document_name, structure, metadata):
        """
        Generate a professional PowerPoint presentation
        
        Args:
            content: Generated content from AI
            document_name: Name of the document
            structure: Document structure from PM intelligence
            metadata: Document metadata
        
        Returns:
            Path to generated PowerPoint file
        """
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = PptxInches(10)
            prs.slide_height = PptxInches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = document_name
            subtitle.text = f"{metadata.get('methodology', 'PMI')} Methodology\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Parse content and add slides
            self._add_powerpoint_slides(prs, content, structure, metadata)
            
            # Save presentation
            filename = f"{document_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(self.temp_dir, filename)
            prs.save(filepath)
            
            logger.info(f"PowerPoint document generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint document: {e}")
            raise
    
    # Helper methods
    
    def _setup_word_styles(self, doc):
        """Set up professional Word document styles"""
        styles = doc.styles
        
        # Heading 1
        heading1 = styles['Heading 1']
        heading1.font.name = 'Arial'
        heading1.font.size = Pt(16)
        heading1.font.bold = True
        heading1.font.color.rgb = RGBColor(0, 102, 204)
        
        # Heading 2
        heading2 = styles['Heading 2']
        heading2.font.name = 'Arial'
        heading2.font.size = Pt(14)
        heading2.font.bold = True
        
        # Normal
        normal = styles['Normal']
        normal.font.name = 'Arial'
        normal.font.size = Pt(11)
    
    def _add_word_title_page(self, doc, document_name, metadata):
        """Add professional title page to Word document"""
        # Title
        title = doc.add_heading(document_name, level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add spacing
        doc.add_paragraph()
        doc.add_paragraph()
        
        # Metadata table
        table = doc.add_table(rows=5, cols=2)
        table.style = 'Light Grid Accent 1'
        
        table.cell(0, 0).text = 'Methodology:'
        table.cell(0, 1).text = metadata.get('methodology', 'N/A')
        
        table.cell(1, 0).text = 'PMBOK Knowledge Area:'
        table.cell(1, 1).text = metadata.get('pmbok_knowledge_area', 'N/A')
        
        table.cell(2, 0).text = 'Process Group:'
        table.cell(2, 1).text = metadata.get('pmbok_process_group', 'N/A')
        
        table.cell(3, 0).text = 'Generated:'
        table.cell(3, 1).text = datetime.now().strftime('%B %d, %Y at %H:%M')
        
        table.cell(4, 0).text = 'Platform:'
        table.cell(4, 1).text = 'PMBlueprints - Professional PM Templates'
    
    def _add_word_content(self, doc, content, structure):
        """Parse AI content and add to Word document"""
        # Split content into sections
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Detect headings (lines starting with # or all caps)
            if line.startswith('#'):
                level = line.count('#')
                text = line.replace('#', '').strip()
                doc.add_heading(text, level=min(level, 3))
            elif line.isupper() and len(line) < 100:
                doc.add_heading(line.title(), level=2)
            elif line.startswith('-') or line.startswith('•'):
                # Bullet point
                doc.add_paragraph(line.lstrip('-•').strip(), style='List Bullet')
            elif line.startswith(tuple(str(i) + '.' for i in range(1, 10))):
                # Numbered list
                doc.add_paragraph(line.split('.', 1)[1].strip(), style='List Number')
            else:
                # Regular paragraph
                doc.add_paragraph(line)
    
    def _add_word_footer(self, doc, metadata):
        """Add footer to Word document"""
        section = doc.sections[0]
        footer = section.footer
        footer_para = footer.paragraphs[0]
        footer_para.text = f"PMBlueprints | {metadata.get('methodology', 'PMI')} | Page "
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    def _add_excel_sample_data(self, ws, start_row, columns, content, metadata):
        """Add sample data rows to Excel based on AI content"""
        # Parse content for data
        sample_data = self._parse_content_for_excel(content, columns, metadata)
        
        for row_idx, row_data in enumerate(sample_data, start=start_row):
            for col_idx, value in enumerate(row_data, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
    
    def _parse_content_for_excel(self, content, columns, metadata):
        """Parse AI content to extract data for Excel rows"""
        # Generate sample rows based on document type
        methodology = metadata.get('methodology', 'waterfall').lower()
        
        sample_data = []
        for i in range(1, 11):  # Generate 10 sample rows
            row = []
            for col in columns:
                col_lower = col.lower()
                if 'id' in col_lower:
                    row.append(f"{i:03d}")
                elif 'description' in col_lower or 'name' in col_lower:
                    row.append(f"Sample {col} {i}")
                elif 'owner' in col_lower or 'assigned' in col_lower:
                    row.append(f"Team Member {i % 3 + 1}")
                elif 'status' in col_lower:
                    statuses = ['Not Started', 'In Progress', 'Completed', 'On Hold']
                    row.append(statuses[i % len(statuses)])
                elif 'priority' in col_lower:
                    priorities = ['High', 'Medium', 'Low']
                    row.append(priorities[i % len(priorities)])
                elif 'date' in col_lower:
                    row.append(datetime.now().strftime('%Y-%m-%d'))
                else:
                    row.append(f"Data {i}")
            sample_data.append(row)
        
        return sample_data
    
    def _add_powerpoint_slides(self, prs, content, structure, metadata):
        """Add content slides to PowerPoint presentation"""
        # Agenda slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Agenda'
        
        tf = body_shape.text_frame
        tf.text = 'Overview'
        
        p = tf.add_paragraph()
        p.text = 'Key Components'
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = 'Implementation Plan'
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = 'Next Steps'
        p.level = 0
        
        # Content slides - parse AI content
        sections = self._parse_content_for_powerpoint(content)
        
        for section_title, section_content in sections:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = section_title
            
            tf = body_shape.text_frame
            tf.text = section_content[0] if section_content else ''
            
            for item in section_content[1:5]:  # Limit to 4 bullets per slide
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
        
        # Closing slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = 'Questions?'
        subtitle.text = f"Generated by PMBlueprints\n{metadata.get('methodology', 'PMI')} Methodology"
    
    def _parse_content_for_powerpoint(self, content):
        """Parse AI content into sections for PowerPoint slides"""
        sections = []
        current_section = None
        current_items = []
        
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Detect section headers
            if line.startswith('#') or (line.isupper() and len(line) < 100):
                if current_section:
                    sections.append((current_section, current_items))
                current_section = line.replace('#', '').strip().title()
                current_items = []
            elif line.startswith('-') or line.startswith('•'):
                current_items.append(line.lstrip('-•').strip())
            elif current_section and len(current_items) < 5:
                current_items.append(line[:100])  # Limit line length
        
        if current_section:
            sections.append((current_section, current_items))
        
        return sections[:8]  # Limit to 8 content slides


# Global instance
document_generator = DocumentGenerator()

